"""Tests for the PPTX deck exporter (#66)."""

from __future__ import annotations

from pathlib import Path

import pytest

from app.pptx_export import export_deck

pytest.importorskip("pptx")

SPEC = [
    {"title": "Trust in automated alerts", "bullets": ["How people trust alerts."], "notes": "", "draft": False},
    {
        "title": "[draft] control-earns-trust",
        "bullets": ["Trust rises with manual override."],
        "notes": 'U-0002 (S001): "no sé si confiar"',
        "draft": True,
    },
    {"title": "Saturation", "bullets": ["Recommendation: pause"], "notes": "", "draft": False},
]


def test_export_deck_writes_one_slide_per_spec_entry(tmp_path: Path):
    from pptx import Presentation

    out = tmp_path / "deck.pptx"
    export_deck(SPEC, str(out))
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_citations_become_slide_notes(tmp_path: Path):
    from pptx import Presentation

    out = tmp_path / "deck.pptx"
    export_deck(SPEC, str(out))
    prs = Presentation(str(out))
    notes = prs.slides[1].notes_slide.notes_text_frame.text
    assert "U-0002" in notes
    assert "no sé si confiar" in notes


def test_title_carries_draft_marker(tmp_path: Path):
    from pptx import Presentation

    out = tmp_path / "deck.pptx"
    export_deck(SPEC, str(out))
    prs = Presentation(str(out))
    assert prs.slides[1].shapes.title.text.startswith("[draft]")
