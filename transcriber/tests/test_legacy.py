"""Tests for legacy ingestors (#29).

Fixtures are generated programmatically (no binary blobs in the repo). Each
ingestor's output is validated against schema/transcript.schema.json.
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest

from app.legacy import (
    _build_speaker_utterances,
    _collect_label_counts,
    _speaker_turns,
    _strip_running_headers,
    _valid_speaker_names,
    ingest,
    ingest_csv,
    ingest_docx,
    ingest_pdf,
    ingest_pptx,
)

SCHEMA_PATH = Path(__file__).resolve().parents[2] / "schema" / "transcript.schema.json"


def _validate(doc: dict) -> None:
    import jsonschema  # type: ignore

    schema = json.loads(SCHEMA_PATH.read_text())
    jsonschema.validate(doc, schema)


def test_csv_one_utterance_per_row_with_mapping(tmp_path: Path):
    csv_path = tmp_path / "quotes.csv"
    csv_path.write_text("Quote,Participant\n\"I distrust alerts\",P07\n\"It depends\",P08\n")
    doc = ingest_csv(csv_path, text_col="Quote", speaker_col="Participant")
    assert doc["kind"] == "document"
    assert len(doc["utterances"]) == 2
    assert doc["utterances"][0]["text"] == "I distrust alerts"
    assert "[speaker: P07]" in doc["utterances"][0]["annotation"]
    _validate(doc)


def test_csv_missing_column_raises(tmp_path: Path):
    csv_path = tmp_path / "q.csv"
    csv_path.write_text("A,B\n1,2\n")
    with pytest.raises(ValueError):
        ingest_csv(csv_path, text_col="Quote")


def test_docx_one_utterance_per_paragraph_with_section_anchor(tmp_path: Path):
    import docx

    d = docx.Document()
    d.add_heading("Findings", level=1)
    d.add_paragraph("Participants distrust automated alerts.")
    d.add_paragraph("They want a manual override.")
    p = tmp_path / "report.docx"
    d.save(p)

    doc = ingest_docx(p)
    assert len(doc["utterances"]) == 2
    assert doc["utterances"][0]["annotation"] == "[section: Findings]"
    _validate(doc)


def test_pptx_one_utterance_per_slide(tmp_path: Path):
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title in ("Intro", "Key Insight"):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(0, 0, 100, 100).text_frame
        tb.text = title
    p = tmp_path / "deck.pptx"
    prs.save(p)

    doc = ingest_pptx(p)
    assert len(doc["utterances"]) == 2
    assert doc["utterances"][0]["source_page"] == 1
    assert "Intro" in doc["utterances"][0]["text"]
    _validate(doc)


def test_pdf_paragraphs_with_source_page(tmp_path: Path):
    # Build a tiny PDF with reportlab if available; else skip (text-extraction
    # is exercised by the paragraph splitter test below regardless).
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas  # type: ignore

    p = tmp_path / "doc.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(72, 720, "First finding about trust.")
    c.showPage()
    c.drawString(72, 720, "Second finding about control.")
    c.showPage()
    c.save()

    doc = ingest_pdf(p)
    assert doc["kind"] == "document"
    assert len(doc["utterances"]) >= 2
    pages = {u["source_page"] for u in doc["utterances"]}
    assert pages == {1, 2}
    _validate(doc)


def test_dispatch_by_extension(tmp_path: Path):
    csv_path = tmp_path / "q.csv"
    csv_path.write_text("text\nhello\n")
    doc = ingest(csv_path, text_col="text")
    assert doc["kind"] == "document"
    assert doc["utterances"][0]["text"] == "hello"


def test_unsupported_extension_raises(tmp_path: Path):
    with pytest.raises(ValueError):
        ingest(tmp_path / "x.rtf")


# #270: source/author attribution for sourced documents (standpoint w/o speaker).


def test_attribution_attached_via_ingest(tmp_path: Path):
    csv_path = tmp_path / "essay.csv"
    csv_path.write_text("text\nA situated claim about worlds.\n")
    doc = ingest(
        csv_path,
        text_col="text",
        attribution={"author": "Donna Haraway", "year": "2007", "title": "Edges & Ecotones"},
    )
    assert doc["attribution"] == {
        "author": "Donna Haraway",
        "year": "2007",
        "title": "Edges & Ecotones",
    }
    _validate(doc)


def test_attribution_drops_empty_and_unknown_keys(tmp_path: Path):
    csv_path = tmp_path / "q.csv"
    csv_path.write_text("text\nhello\n")
    doc = ingest(
        csv_path,
        text_col="text",
        attribution={"author": "  ", "year": "1990", "speaker": "should be ignored"},
    )
    # Empty author dropped, unknown 'speaker' dropped, year kept.
    assert doc["attribution"] == {"year": "1990"}
    _validate(doc)


def test_no_attribution_when_all_empty_or_absent(tmp_path: Path):
    csv_path = tmp_path / "q.csv"
    csv_path.write_text("text\nhello\n")
    assert "attribution" not in ingest(csv_path, text_col="text")
    assert "attribution" not in ingest(csv_path, text_col="text", attribution={"author": ""})


# v0.1-02 review feedback: auto-detect column name when text_col is None.


def test_csv_autodetects_text_column_when_text_col_is_none(tmp_path: Path):
    """Header has `transcript` (not `text`). Auto-detect finds it second
    in the priority list."""
    csv_path = tmp_path / "otter.csv"
    csv_path.write_text("speaker,transcript\nMod,hello there\nP01,thank you\n")
    doc = ingest_csv(csv_path)  # text_col omitted → auto-detect
    assert doc["provenance"]["text_col_resolved"] == "transcript"
    assert [u["text"] for u in doc["utterances"]] == ["hello there", "thank you"]


def test_csv_autodetect_is_case_insensitive(tmp_path: Path):
    """The header says `Content` (capital C); auto-detect should still match."""
    csv_path = tmp_path / "zoom.csv"
    csv_path.write_text("Speaker,Content\nMod,intro\nP01,follow-up\n")
    doc = ingest_csv(csv_path)
    assert doc["provenance"]["text_col_resolved"] == "Content"
    assert len(doc["utterances"]) == 2


def test_csv_autodetect_falls_back_to_first_column(tmp_path: Path):
    """Header has no candidate from the priority list → fall back to col 0."""
    csv_path = tmp_path / "weird.csv"
    csv_path.write_text("notes,extra\nfirst row,x\nsecond row,y\n")
    doc = ingest_csv(csv_path)
    assert doc["provenance"]["text_col_resolved"] == "notes"
    assert [u["text"] for u in doc["utterances"]] == ["first row", "second row"]


def test_csv_explicit_text_col_overrides_autodetect(tmp_path: Path):
    """When text_col is passed explicitly, auto-detect is skipped."""
    csv_path = tmp_path / "survey.csv"
    csv_path.write_text("text,answer\nshould be ignored,real answer\n")
    doc = ingest_csv(csv_path, text_col="answer")
    assert doc["provenance"]["text_col_resolved"] == "answer"
    assert doc["utterances"][0]["text"] == "real answer"


# #271 seam 2: strip repeated PDF running headers/footers before segmentation.


def test_strip_running_headers_removes_repeated_zone_lines():
    # The header recurs on every page, varying only by trailing page number.
    pages = [
        "Edges and Ecotones: Worlds at UCSC 1\n\nFirst real paragraph of the interview.",
        "Edges and Ecotones: Worlds at UCSC 2\n\nSecond real paragraph continues here.",
        "Edges and Ecotones: Worlds at UCSC 3\n\nThird real paragraph wraps up.",
    ]
    out = _strip_running_headers(pages)
    joined = "\n".join(out)
    assert "Worlds at UCSC" not in joined
    assert "First real paragraph" in joined
    assert "Third real paragraph" in joined


def test_strip_running_headers_keeps_repeated_body_text():
    # An identical sentence in the *body* (outside the top/bottom 3-line zone)
    # must survive even though it recurs across pages.
    pages = [
        "Header A 1\nb\nc\nd\nThe refrain repeats.\ne\nf\ng\nFooter Z 1",
        "Header A 2\nb2\nc2\nd2\nThe refrain repeats.\ne2\nf2\ng2\nFooter Z 2",
    ]
    out = _strip_running_headers(pages)
    joined = "\n".join(out)
    assert joined.count("The refrain repeats.") == 2  # body refrain kept
    assert "Header A" not in joined
    assert "Footer Z" not in joined


def test_strip_running_headers_noop_single_page():
    pages = ["Only one page here 1\n\nBody."]
    assert _strip_running_headers(pages) == pages


def test_strip_running_headers_keeps_body_on_sparse_pages():
    # Review finding: on short pages (<= 2*zone non-empty lines) the top and
    # bottom bands must not overlap and swallow a recurring *body* sentence.
    # Header + recurring middle body line + footer, only 3 lines per page.
    pages = [
        "Worlds at UCSC 1\nWe studied worlds.\nfooter alpha 1",
        "Worlds at UCSC 2\nWe studied worlds.\nfooter alpha 2",
    ]
    out = _strip_running_headers(pages)
    joined = "\n".join(out)
    assert joined.count("We studied worlds.") == 2  # body survives
    assert "Worlds at UCSC" not in joined  # header still stripped
    assert "footer alpha" not in joined  # footer still stripped


# #271 seam 1: split leading "Name:" turn labels into speaker-attributed turns.


def test_speaker_turns_splits_leading_labels():
    valid = {"Reti", "Haraway"}
    turns = _speaker_turns("Reti: So tell me. Haraway: Well, it began in the lab.", valid)
    assert turns == [("Reti", "So tell me."), ("Haraway", "Well, it began in the lab.")]


def test_speaker_turns_unattributed_when_no_leading_label():
    turns = _speaker_turns("A plain paragraph with Haraway: midway through.", {"Haraway"})
    assert turns == [(None, "A plain paragraph with Haraway: midway through.")]


def test_collect_label_counts_filters_stopwords_keeps_raw_counts():
    paras = [
        "Reti: question one",
        "Haraway: answer one",
        "Reti: question two",
        "Abstract: this is a section heading",
        "Note: an aside",
    ]
    counts = _collect_label_counts(paras)
    assert counts["Reti"] == 2
    assert counts["Haraway"] == 1  # raw count; promotion gate is tested separately
    assert "Abstract" not in counts  # stopword
    assert "Note" not in counts  # stopword


def test_valid_speaker_names_excludes_singletons():
    # The >=2 recurrence gate is the heuristic's central safety rail: a label
    # seen once stays prose. Mutating the threshold to >=1 must fail here.
    paras = [
        "Reti: question one",
        "Haraway: answer one",
        "Reti: question two",
        "Smith: a one-off aside that is not a recurring speaker",
    ]
    assert _valid_speaker_names(paras) == {"Reti"}


def test_speaker_turns_handles_adjacent_labels():
    # Review finding: a second label immediately following the first must not be
    # cannibalized by the first match's separator. Reti has no body, Haraway does.
    turns = _speaker_turns("Reti: Haraway: the actual content here", {"Reti", "Haraway"})
    assert turns == [("Haraway", "the actual content here")]


def test_pdf_singleton_label_stays_prose(tmp_path: Path):
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas  # type: ignore

    p = tmp_path / "mostly-prose.pdf"
    c = canvas.Canvas(str(p))
    # Only "Reti" recurs; "Smith" appears once and must NOT become a speaker.
    lines = ["Reti: opening question", "Reti: a follow-up question", "Smith: a single stray label"]
    for line in lines:
        c.drawString(72, 700, line)
        c.showPage()
    c.save()

    doc = ingest_pdf(p)
    names = {s.get("name") for s in doc["speakers"]}
    assert "Reti" in names
    assert "Smith" not in names  # one-off label not promoted
    # The Smith line survives verbatim somewhere as document-speaker prose.
    assert any("Smith: a single stray label" in u["text"] for u in doc["utterances"])


def test_build_speaker_utterances_assigns_ids_and_annotates():
    turns = [
        ("Reti", "question one", 1),
        ("Haraway", "answer one", 1),
        (None, "stray unlabeled prose", 2),
        ("Reti", "question two", 2),
    ]
    speakers, utts = _build_speaker_utterances(turns)
    by_name = {s.get("name"): s["id"] for s in speakers}
    assert by_name["document"] == "S1"  # reserved for unlabeled
    assert utts[0]["speaker_id"] == by_name["Reti"]
    assert utts[0]["speaker_id"] != by_name["Haraway"]
    assert utts[3]["speaker_id"] == by_name["Reti"]  # stable across turns
    assert utts[2]["speaker_id"] == "S1"  # unlabeled → document
    assert utts[2].get("annotation") is None
    assert utts[0]["annotation"] == "[speaker-from-label: Reti]"
    assert all(s["type"] in ("moderator", "participant", "other") for s in speakers)


def test_pdf_strips_headers_and_attributes_speakers(tmp_path: Path):
    pytest.importorskip("reportlab")
    from reportlab.pdfgen import canvas  # type: ignore

    p = tmp_path / "interview.pdf"
    c = canvas.Canvas(str(p))
    # Each speaker recurs (≥2 turns) so both clear the recurrence gate.
    dialogue = [
        ("Reti", "How did the program begin?"),
        ("Haraway", "It began with a question about worlds."),
        ("Reti", "And what came next?"),
        ("Haraway", "Then the worlds multiplied."),
    ]
    for page_no, (name, line) in enumerate(dialogue, start=1):
        c.drawString(72, 760, f"Edges and Ecotones: Worlds at UCSC {page_no}")  # running header
        c.drawString(72, 700, f"{name}: {line}")
        c.showPage()
    c.save()

    doc = ingest_pdf(p)
    full = " ".join(u["text"] for u in doc["utterances"])
    assert "Worlds at UCSC" not in full  # running header stripped
    names = {s.get("name") for s in doc["speakers"]}
    assert {"Reti", "Haraway"} <= names  # both speakers recovered
    speaker_ids = {u["speaker_id"] for u in doc["utterances"]}
    assert len(speaker_ids) >= 2  # not all under one speaker
    _validate(doc)
