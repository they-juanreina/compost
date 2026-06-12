"""Legacy asset ingestors (#29).

Normalize PDF / DOCX / PPTX / CSV into a transcript-shaped JSON with
kind="document": one utterance per paragraph (PDF/DOCX), per slide (PPTX),
or per row (CSV). Output validates against schema/transcript.schema.json
(kind="document", modality=["document"]).

Heavy parsers (pdfplumber, python-docx, python-pptx) are imported lazily so
the module loads without the `legacy` extra; each ingestor raises a clear
error if its dependency is missing.
"""

from __future__ import annotations

import csv
import os
import re
from collections import Counter
from pathlib import Path
from typing import Any

INGESTOR_VERSION = "compost-legacy@0.1.0"
DOC_SPEAKER = {"id": "S1", "name": "document", "type": "other"}


def _base(session_id: str, source: str, language: str = "und") -> dict[str, Any]:
    return {
        "schema_version": "1.0",
        "kind": "document",
        "session_id": session_id,
        "source": source,
        "language": language,
        "duration_ms": 0,
        "modality": ["document"],
        "speakers": [dict(DOC_SPEAKER)],
        "utterances": [],
        "provenance": {"transcriber": INGESTOR_VERSION},
    }


def _utt(idx: int, text: str, source_page: int | None = None, annotation: str | None = None) -> dict[str, Any]:
    u: dict[str, Any] = {
        "id": f"U-{idx:04d}",
        "speaker_id": DOC_SPEAKER["id"],
        "turn": idx,
        "start_ms": 0,
        "end_ms": 0,
        "text": text,
    }
    if source_page is not None:
        u["source_page"] = source_page
    if annotation is not None:
        u["annotation"] = annotation
    return u


def _session_id(path: str | Path) -> str:
    stem = Path(path).stem
    safe = "".join(c if c.isalnum() or c in "-_" else "-" for c in stem)
    return f"DOC-{safe}"[:64]


# ---------------------------------------------------------------- CSV / XLSX

# Auto-detect priority for the "text" column. First case-insensitive match
# in the source's header wins. Falls back to the first column.
TEXT_COL_CANDIDATES = (
    "text",
    "transcript",
    "content",
    "utterance",
    "quote",
    "message",
    "body",
)


def _auto_text_col(fieldnames: list[str]) -> str:
    """Pick the most-likely text column from a header. Case-insensitive match
    against TEXT_COL_CANDIDATES, then a first-column fallback."""
    lower = {f.lower(): f for f in fieldnames}
    for candidate in TEXT_COL_CANDIDATES:
        if candidate in lower:
            return lower[candidate]
    return fieldnames[0]


def ingest_csv(
    path: str | Path,
    text_col: str | None = None,
    speaker_col: str | None = None,
) -> dict[str, Any]:
    """One utterance per row.

    `text_col=None` triggers auto-detect: text → transcript → content →
    utterance → quote → message → body (case-insensitive). Falls back to
    the first column. The resolved column is recorded on the output's
    `provenance.text_col_resolved` for caller visibility.
    """
    path = str(path)
    doc = _base(_session_id(path), path)
    with open(path, newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        if reader.fieldnames is None:
            raise ValueError(f"CSV has no header row: {path}")
        fields = list(reader.fieldnames)
        resolved = text_col if text_col is not None else _auto_text_col(fields)
        if resolved not in fields:
            raise ValueError(f"CSV has no column '{resolved}' (columns: {fields})")
        doc["provenance"]["text_col_resolved"] = resolved
        idx = 1
        for row in reader:
            text = (row.get(resolved) or "").strip()
            if not text:
                continue
            ann = None
            if speaker_col is not None and row.get(speaker_col):
                ann = f"[speaker: {row[speaker_col]}]"
            doc["utterances"].append(_utt(idx, text, source_page=idx, annotation=ann))
            idx += 1
    return doc


# ---------------------------------------------------------------- DOCX


def ingest_docx(path: str | Path) -> dict[str, Any]:
    try:
        import docx  # type: ignore
    except ImportError as e:
        raise RuntimeError("python-docx not installed (pip install -e '.[legacy]')") from e

    path = str(path)
    doc = _base(_session_id(path), path)
    d = docx.Document(path)
    idx = 1
    current_heading: str | None = None
    for para in d.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            current_heading = text
            # headings preserved as section anchors via annotation on the next utterances
            continue
        ann = f"[section: {current_heading}]" if current_heading else None
        doc["utterances"].append(_utt(idx, text, annotation=ann))
        idx += 1
    return doc


# ---------------------------------------------------------------- PPTX


def ingest_pptx(path: str | Path, thumbnails_dir: str | Path | None = None) -> dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as e:
        raise RuntimeError("python-pptx not installed (pip install -e '.[legacy]')") from e

    path = str(path)
    doc = _base(_session_id(path), path)
    prs = Presentation(path)
    idx = 1
    for slide_no, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in p.runs).strip()
                    if line:
                        parts.append(line)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(f"(notes) {notes}")
        text = "\n".join(parts)
        if text:
            doc["utterances"].append(_utt(idx, text, source_page=slide_no))
            idx += 1
    # Thumbnail rendering requires LibreOffice/unoconv (not bundled); skipped
    # gracefully. The slide text above is the load-bearing evidence.
    if thumbnails_dir is not None:
        os.makedirs(thumbnails_dir, exist_ok=True)
    return doc


# ---------------------------------------------------------------- PDF


def ingest_pdf(path: str | Path) -> dict[str, Any]:
    try:
        import pdfplumber  # type: ignore
    except ImportError as e:
        raise RuntimeError("pdfplumber not installed (pip install -e '.[legacy]')") from e

    path = str(path)
    doc = _base(_session_id(path), path)
    raw_pages: list[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            # OCR fallback for scanned pages (no extractable text) requires
            # pytesseract + the page raster; attempted best-effort.
            if not text.strip():
                text = _ocr_page(page)
            raw_pages.append(text)

    # Seam 2 (#271): each page's running header/footer ("…Worlds at UCSC N")
    # is otherwise absorbed into the first/last paragraph of every page. Strip
    # lines that recur in the header/footer zone across pages before segmenting.
    pages = _strip_running_headers(raw_pages)

    # Seam 1 (#271): two-speaker oral histories land under a single speaker with
    # "Reti:"/"Haraway:" labels inside the paragraph text. Promote a label to a
    # speaker only when it recurs (≥2 turns), so one-off section labels
    # ("Abstract:", "Note:") are left as prose. Best-effort, researcher-overridable.
    paragraphs = [(para, page_no) for page_no, text in enumerate(pages, start=1) for para in _paragraphs(text)]
    valid_names = _valid_speaker_names([p for p, _ in paragraphs])

    turns: list[tuple[str | None, str, int]] = []
    for para, page_no in paragraphs:
        for name, body in _speaker_turns(para, valid_names):
            turns.append((name, body, page_no))

    doc["speakers"], doc["utterances"] = _build_speaker_utterances(turns)
    return doc


# ---------------------------------------------------------------- header/footer

_RUNNING_ZONE = 3  # lines from the top and bottom of a page treated as header/footer
_DIGITS_RE = re.compile(r"\d+")
_WS_RE = re.compile(r"\s+")


def _normalize_running(line: str) -> str:
    """Signature for a candidate running-header line: whitespace collapsed and
    digit runs replaced by `#`, so the page number that varies per page does
    not defeat the cross-page match."""
    return _DIGITS_RE.sub("#", _WS_RE.sub(" ", line).strip()).casefold()


def _zone_indices(lines: list[str], zone: int) -> set[int]:
    """Line indices in the top/bottom header/footer band of one page.

    The per-side band is shrunk on short pages so it never covers an *interior*
    line — one with content both above and below it. That is the guarantee: a
    sentence legitimately repeated in body text (which by definition sits
    between other lines) is never mistaken for a running header. A 3-line page
    keeps its middle line; a 1–2-line page has no interior to protect, so its
    top/bottom line can still be stripped if it recurs.
    """
    nonempty = [i for i, ln in enumerate(lines) if ln.strip()]
    if not nonempty:
        return set()
    # eff <= (N-1)//2 guarantees an interior line survives whenever N >= 3.
    eff = min(zone, max(1, (len(nonempty) - 1) // 2))
    return set(nonempty[:eff]) | set(nonempty[-eff:])


def _strip_running_headers(pages: list[str], zone: int = _RUNNING_ZONE) -> list[str]:
    """Remove lines that recur in the top/bottom `zone` of a majority of pages.

    Only the header/footer band (see `_zone_indices`) is considered, so a
    sentence legitimately repeated in body text is never stripped. A no-op for
    single-page inputs.
    """
    if len(pages) < 2:
        return pages
    page_lines = [p.splitlines() for p in pages]

    counts: Counter[str] = Counter()
    for lines in page_lines:
        for sig in {_normalize_running(lines[i]) for i in _zone_indices(lines, zone)}:
            if sig:
                counts[sig] += 1

    threshold = max(2, (len(pages) + 1) // 2)
    running = {sig for sig, c in counts.items() if c >= threshold}
    if not running:
        return pages

    out: list[str] = []
    for lines in page_lines:
        zone_idx = _zone_indices(lines, zone)
        kept = [ln for i, ln in enumerate(lines) if not (i in zone_idx and _normalize_running(ln) in running)]
        out.append("\n".join(kept))
    return out


# ---------------------------------------------------------------- speaker labels

# A leading turn label: an initial-capitalized name (internal apostrophes/hyphens
# allowed, e.g. "O'Brien", "Smith-Jones") followed by a colon and whitespace.
# Start (`^`/lookbehind) and trailing-space (lookahead) are zero-width so that
# `finditer` — which is non-overlapping — still matches a second label that
# immediately follows the first ("Reti: Haraway: …") instead of the first
# match's trailing space cannibalizing the second's separator.
_TURN_LABEL_RE = re.compile(r"(?:^|(?<=\s))([A-Z][A-Za-z.'\-]{1,20}):(?=\s)")

# Labels that look like names but introduce a section/field, not a speaker turn.
_LABEL_STOPWORDS = frozenset({
    "abstract", "note", "notes", "see", "figure", "fig", "table", "page",
    "vol", "no", "example", "source", "sources", "keywords", "introduction",
    "conclusion", "summary", "ibid", "http", "https", "www",
})


def _is_plausible_label(name: str) -> bool:
    return len(name) >= 2 and name.casefold() not in _LABEL_STOPWORDS


def _collect_label_counts(paragraphs: Any) -> Counter[str]:
    counts: Counter[str] = Counter()
    for para in paragraphs:
        for m in _TURN_LABEL_RE.finditer(para):
            if _is_plausible_label(m.group(1)):
                counts[m.group(1)] += 1
    return counts


# A label is promoted to a speaker only if it recurs. A one-off label — a stray
# "Smith:" or a section heading the stopword list missed — appears once and so
# stays as unattributed prose. This >=2 threshold is the heuristic's safety rail
# against minting phantom speakers from incidental colon-prefixed words.
_SPEAKER_RECURRENCE_MIN = 2


def _valid_speaker_names(paragraphs: list[str]) -> set[str]:
    counts = _collect_label_counts(paragraphs)
    return {name for name, count in counts.items() if count >= _SPEAKER_RECURRENCE_MIN}


def _speaker_turns(text: str, valid_names: set[str]) -> list[tuple[str | None, str]]:
    """Split a paragraph into (speaker|None, body) turns on recurring labels.

    Conservative: only splits when the paragraph *begins* with a recognized
    label (honoring the issue's "leading Name: turn label"); otherwise the
    paragraph is returned whole as an unattributed turn.
    """
    matches = [m for m in _TURN_LABEL_RE.finditer(text) if m.group(1) in valid_names]
    if not matches or matches[0].start() != 0:
        return [(None, text)]
    turns: list[tuple[str | None, str]] = []
    for i, m in enumerate(matches):
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[m.end():end].strip()
        if body:
            turns.append((m.group(1), body))
    return turns


def _build_speaker_utterances(
    turns: list[tuple[str | None, str, int]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Map (name|None, text, page) turns to schema speakers + utterances.

    Unattributed turns share the reserved `document` speaker (S1); each distinct
    label gets the next Sn id. Speaker-from-label utterances are annotated so a
    researcher can see — and override — the heuristic attribution.
    """
    has_unlabeled = any(name is None for name, _, _ in turns)
    speakers: list[dict[str, Any]] = []
    name_to_id: dict[str, str] = {}
    next_n = 1
    if has_unlabeled or not turns:
        speakers.append(dict(DOC_SPEAKER))  # S1 = document
        next_n = 2

    utterances: list[dict[str, Any]] = []
    for idx, (name, body, page_no) in enumerate(turns, start=1):
        if name is None:
            speaker_id = DOC_SPEAKER["id"]
            annotation = None
        else:
            if name not in name_to_id:
                speaker_id = f"S{next_n}"
                next_n += 1
                name_to_id[name] = speaker_id
                speakers.append({"id": speaker_id, "name": name, "type": "other"})
            speaker_id = name_to_id[name]
            annotation = f"[speaker-from-label: {name}]"
        u = _utt(idx, body, source_page=page_no, annotation=annotation)
        u["speaker_id"] = speaker_id
        utterances.append(u)

    if not speakers:
        speakers.append(dict(DOC_SPEAKER))
    return speakers, utterances


def _paragraphs(text: str) -> list[str]:
    out: list[str] = []
    for block in text.split("\n\n"):
        cleaned = " ".join(line.strip() for line in block.splitlines() if line.strip())
        if cleaned:
            out.append(cleaned)
    return out


def _ocr_page(page: Any) -> str:  # pragma: no cover - needs tesseract + a raster
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore  # noqa: F401
    except ImportError:
        return ""
    try:
        im = page.to_image(resolution=200).original
        return pytesseract.image_to_string(im)
    except Exception:
        return ""


# ---------------------------------------------------------------- Markdown / Text


def ingest_text(path: str | Path) -> dict[str, Any]:
    """Read a plain-text or Markdown file and split into paragraph utterances.

    Both `.txt` (Otter / Zoom exports) and `.md` land here. Top-level
    headings are recorded as section annotations on subsequent utterances,
    mirroring the docx behavior.
    """
    path = str(path)
    doc = _base(_session_id(path), path)
    with open(path, encoding="utf-8") as f:
        body = f.read()

    current_heading: str | None = None
    idx = 1
    for para in _paragraphs(body):
        # Markdown ATX heading line (`# ` through `###### `) → record as section
        # anchor, skip the utterance. Setext (==== / ---- underline) not yet
        # supported — rare in mod-era markdown.
        if para.startswith(("# ", "## ", "### ", "#### ", "##### ", "###### ")):
            current_heading = para.lstrip("# ").strip()
            continue
        ann = f"[section: {current_heading}]" if current_heading else None
        doc["utterances"].append(_utt(idx, para, annotation=ann))
        idx += 1
    return doc


# ---------------------------------------------------------------- XLSX


def ingest_xlsx(
    path: str | Path,
    text_col: str | None = None,
    speaker_col: str | None = None,
    sheet: str | None = None,
) -> dict[str, Any]:
    """One utterance per row of a spreadsheet.

    `text_col=None` triggers the same auto-detect as `ingest_csv`. The
    resolved column lands on `provenance.text_col_resolved`. Use `sheet`
    to pick a non-default tab.
    """
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as e:
        raise RuntimeError("openpyxl not installed (pip install -e '.[legacy]')") from e

    path = str(path)
    doc = _base(_session_id(path), path)
    wb = load_workbook(path, read_only=True, data_only=True)
    ws = wb[sheet] if sheet is not None else wb.active
    if ws is None:
        raise ValueError(f"XLSX has no worksheets: {path}")

    rows = ws.iter_rows(values_only=True)
    header_row = next(rows, None)
    if header_row is None:
        return doc  # empty sheet
    header = [str(c) if c is not None else "" for c in header_row]
    resolved = text_col if text_col is not None else _auto_text_col(header)
    if resolved not in header:
        raise ValueError(f"XLSX has no column '{resolved}' (columns: {header})")
    doc["provenance"]["text_col_resolved"] = resolved
    text_idx = header.index(resolved)
    speaker_idx = header.index(speaker_col) if speaker_col in header else -1

    utt_idx = 1
    # Track rows where the text column is empty but the row has other data —
    # a strong proxy for "Excel never evaluated this formula so openpyxl
    # returned None". Researchers seeing this should open the file in Excel
    # once or pre-export to CSV.
    rows_with_data_but_empty_text = 0
    for row in rows:
        if row is None:
            continue
        cell = row[text_idx] if text_idx < len(row) else None
        text = str(cell).strip() if cell is not None else ""
        if not text:
            # Row has data elsewhere → likely an un-evaluated formula in the text column.
            if any(c is not None and str(c).strip() for c in row):
                rows_with_data_but_empty_text += 1
            continue
        ann = None
        if speaker_idx >= 0 and speaker_idx < len(row) and row[speaker_idx] is not None:
            ann = f"[speaker: {row[speaker_idx]}]"
        doc["utterances"].append(_utt(utt_idx, text, source_page=utt_idx, annotation=ann))
        utt_idx += 1
    if rows_with_data_but_empty_text > 0:
        doc["provenance"]["xlsx_rows_skipped_empty_text"] = rows_with_data_but_empty_text
    return doc


# Source/author attribution fields (#270). A sourced document (published
# interview, theory text, archival doc) carries the provenance of the *text* —
# the author's standpoint — since there is no diarized participant to attach it
# to. Kept to the schema's `attribution` shape; extra/empty keys are dropped.
ATTRIBUTION_FIELDS = ("author", "title", "year", "url")
# Structured (CSL-JSON-flavored) citation sub-fields (#270). `editors` is a list;
# the rest are scalars. `raw` is the free-form fallback.
CITATION_FIELDS = ("type", "container_title", "pages", "doi", "raw")


def _citation(raw: Any) -> dict[str, Any] | None:
    """Normalize the citation into the structured shape. A bare string becomes
    `{raw: <string>}` (back-compat); a dict is filtered to the known fields,
    with `editors` kept as a non-empty list of strings."""
    if isinstance(raw, str):
        return {"raw": raw.strip()} if raw.strip() else None
    if not isinstance(raw, dict):
        return None
    out: dict[str, Any] = {
        k: str(raw[k]).strip() for k in CITATION_FIELDS if str(raw.get(k, "")).strip()
    }
    editors = raw.get("editors")
    if isinstance(editors, list):
        eds = [str(e).strip() for e in editors if str(e).strip()]
        if eds:
            out["editors"] = eds
    return out or None


def _attribution(raw: Any) -> dict[str, Any] | None:
    """Pick the known attribution fields with non-empty values, or None.
    Best-effort and researcher-overridable: an unknown key is ignored rather
    than failing the ingest, and an all-empty mapping yields no field at all.
    `citation` is normalized to the structured shape (a bare string folds to
    `{raw: ...}`)."""
    if not isinstance(raw, dict):
        return None
    out: dict[str, Any] = {
        k: str(raw[k]).strip() for k in ATTRIBUTION_FIELDS if str(raw.get(k, "")).strip()
    }
    citation = _citation(raw.get("citation"))
    if citation is not None:
        out["citation"] = citation
    return out or None


def ingest(path: str | Path, **kwargs: Any) -> dict[str, Any]:
    """Dispatch by extension. `text_col=None` (the default) triggers
    auto-detect on CSV/XLSX inputs. `attribution={author, title, year,
    citation, url}` (any subset) records the source text's provenance (#270)."""
    ext = Path(path).suffix.lower()
    if ext == ".csv":
        doc = ingest_csv(
            path,
            text_col=kwargs.get("text_col"),
            speaker_col=kwargs.get("speaker_col"),
        )
    elif ext == ".docx":
        doc = ingest_docx(path)
    elif ext == ".pptx":
        doc = ingest_pptx(path, thumbnails_dir=kwargs.get("thumbnails_dir"))
    elif ext == ".pdf":
        doc = ingest_pdf(path)
    elif ext in (".txt", ".md", ".markdown"):
        doc = ingest_text(path)
    elif ext == ".xlsx":
        doc = ingest_xlsx(
            path,
            text_col=kwargs.get("text_col"),
            speaker_col=kwargs.get("speaker_col"),
            sheet=kwargs.get("sheet"),
        )
    else:
        raise ValueError(f"Unsupported legacy asset: {ext}")

    attribution = _attribution(kwargs.get("attribution"))
    if attribution is not None:
        doc["attribution"] = attribution
    return doc
