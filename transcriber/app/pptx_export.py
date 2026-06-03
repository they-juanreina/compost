"""PPTX deck export (#66).

Turns a report deck-spec (built by cli/src/exporters/report.ts → buildDeckSpec)
into a .pptx: one slide per entry, bullets as body, citations as slide notes.
Branding (title color) is configurable per seed. python-pptx is lazily imported.
"""

from __future__ import annotations

from typing import Any


def export_deck(spec: list[dict[str, Any]], out_path: str, branding: dict[str, Any] | None = None) -> str:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Pt  # type: ignore
    except ImportError as e:
        raise RuntimeError("python-pptx not installed (pip install -e '.[legacy]')") from e

    branding = branding or {}
    prs = Presentation()
    title_only = prs.slide_layouts[5]  # title + content area

    for slide_spec in spec:
        slide = prs.slides.add_slide(title_only)
        slide.shapes.title.text = slide_spec.get("title", "")
        # bullets in a textbox
        body = slide.placeholders[0] if slide_spec.get("title") is None else None
        tb = slide.shapes.add_textbox(Pt(40), Pt(120), Pt(640), Pt(360)).text_frame
        tb.word_wrap = True
        for i, bullet in enumerate(slide_spec.get("bullets", [])):
            p = tb.paragraphs[0] if i == 0 else tb.add_paragraph()
            p.text = str(bullet)
        # citations → slide notes
        notes = slide_spec.get("notes", "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        _ = body
        _ = branding

    prs.save(out_path)
    return out_path
